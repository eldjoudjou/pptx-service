"""
PPTX Service — Micro-service de manipulation PowerPoint (Mode XML Pur)

Architecture :
  1. Reçoit une demande utilisateur + référence à un PPTX (ou demande de création)
  2. Inspecte le fichier (structure, contenu)
  3. Appelle un LLM pour planifier les modifications (JSON)
  4. Pour chaque slide à modifier, le LLM retourne le XML modifié directement
  5. Valide le XML, repackage le PPTX
  6. Upload dans la collection SiaGPT

Sécurité : AUCUN exec() — le LLM retourne du XML, pas du code.
"""

import asyncio
import io
import json
import logging
import os
import re
import tempfile
import uuid
from pathlib import Path

import httpx
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
from fastapi.responses import JSONResponse, StreamingResponse
from pptx import Presentation
from pptx.util import Inches
from lxml import etree

import pptx_tools
import pptx_validate

logger = logging.getLogger("pptx-service")
logging.basicConfig(level=logging.INFO)

# ============================================================
# Configuration
# ============================================================

app = FastAPI(title="PPTX Service", version="1.0.0")

# CORS — permettre les appels depuis Langflow/SiaGPT
from fastapi.middleware.cors import CORSMiddleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# LLM API — SiaGPT /plain_llm endpoint
# Format : POST /chat/plain_llm { systemPrompt, query, llm, temperature }
# Modèles disponibles : claude-4.5-sonnet, claude-4.5-haiku, claude-4-sonnet,
#   gpt-5, gpt-5-mini, gpt-4o, gpt-4.1, o3, o4-mini,
#   gemini-2.5-pro, gemini-3-pro, mistral-large-2, etc.
LLM_API_URL = os.environ.get("LLM_API_URL", "https://backend.siagpt.ai/chat/plain_llm")
LLM_API_KEY = os.environ.get("LLM_API_KEY", "")
LLM_MODEL = os.environ.get("LLM_MODEL", "claude-4.5-sonnet")

# SiaGPT Medias API — stockage des fichiers dans la collection
SIAGPT_MEDIAS_URL = os.environ.get("SIAGPT_MEDIAS_URL", "https://backend.siagpt.ai/medias")
SIAGPT_COLLECTION_ID = os.environ.get("SIAGPT_COLLECTION_ID", "")  # UUID de la collection cible

SYSTEM_PROMPT_PATH = os.environ.get("SYSTEM_PROMPT_PATH", "/app/system_prompt.md")
STYLE_CONFIG_PATH = os.environ.get("STYLE_CONFIG_PATH", "/app/sia_config.md")

# Retry
MAX_RETRIES = int(os.environ.get("MAX_RETRIES", "4"))

# ============================================================
# Initialisation
# ============================================================

# Dossier temporaire de travail
Path("/tmp/pptx-work").mkdir(parents=True, exist_ok=True)


# ============================================================
# Fonctions utilitaires — Stockage SiaGPT Medias
# ============================================================

async def save_to_siagpt_medias(data: bytes, filename: str, auth_token: str) -> dict:
    """
    Upload un fichier dans la collection SiaGPT via POST /medias/.
    Retourne les infos du media créé (uuid, name, versions...).
    """
    media_metadata = json.dumps({"collectionId": SIAGPT_COLLECTION_ID})

    async with httpx.AsyncClient(timeout=60.0) as client:
        response = await client.post(
            f"{SIAGPT_MEDIAS_URL}/",
            files={"file": (filename, data, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            data={"media_metadata": media_metadata},
            headers={"Authorization": f"Bearer {auth_token}"},
        )
        response.raise_for_status()
        return response.json()


async def download_from_siagpt_medias(file_uuid: str, auth_token: str) -> tuple[bytes, str]:
    """
    Télécharge un fichier depuis la collection SiaGPT via GET /medias/{uuid}/download.
    Retourne (bytes, filename).
    """
    async with httpx.AsyncClient(timeout=60.0, follow_redirects=True) as client:
        # D'abord récupérer les métadonnées pour le nom du fichier
        meta_response = await client.get(
            f"{SIAGPT_MEDIAS_URL}/{file_uuid}",
            headers={"Authorization": f"Bearer {auth_token}"},
        )
        meta_response.raise_for_status()
        meta = meta_response.json()
        filename = meta.get("name", f"{file_uuid}.pptx")

        # Télécharger le fichier
        dl_response = await client.get(
            f"{SIAGPT_MEDIAS_URL}/{file_uuid}/download",
            headers={"Authorization": f"Bearer {auth_token}"},
        )
        dl_response.raise_for_status()
        return dl_response.content, filename


def load_system_prompt() -> str:
    """
    Charge le system prompt + la config style.
    Le prompt est générique (XML, phases, règles).
    La config est interchangeable (couleurs, polices, layouts).
    """
    try:
        prompt = Path(SYSTEM_PROMPT_PATH).read_text(encoding="utf-8")
    except FileNotFoundError:
        prompt = "Tu es un expert en manipulation PowerPoint via XML. Retourne uniquement du XML."

    try:
        config = Path(STYLE_CONFIG_PATH).read_text(encoding="utf-8")
        prompt += "\n\n---\n\n" + config
        logger.info(f"Config style chargée depuis {STYLE_CONFIG_PATH}")
    except FileNotFoundError:
        logger.warning(f"Config style non trouvée ({STYLE_CONFIG_PATH}) — utilisation du prompt seul")

    return prompt

SYSTEM_PROMPT = load_system_prompt()

# ============================================================
# Inspection PPTX
# ============================================================

def inspect_pptx_structure(pptx_bytes: bytes) -> str:
    """Inspecte la structure complète d'un PPTX, retourne du JSON."""
    prs = Presentation(io.BytesIO(pptx_bytes))

    structure = {
        "slide_width_emu": str(prs.slide_width),
        "slide_height_emu": str(prs.slide_height),
        "slide_count": len(prs.slides),
        "slide_layouts": [],
        "slides": [],
    }

    # Layouts disponibles
    for i, layout in enumerate(prs.slide_layouts):
        structure["slide_layouts"].append({"index": i, "name": layout.name})

    # Contenu de chaque slide
    for i, slide in enumerate(prs.slides):
        slide_info = {
            "index": i,
            "layout": slide.slide_layout.name,
            "shapes": [],
        }
        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left_emu": str(shape.left),
                "top_emu": str(shape.top),
                "width_emu": str(shape.width),
                "height_emu": str(shape.height),
            }
            if shape.has_text_frame:
                shape_info["text"] = shape.text_frame.text[:500]  # Tronquer si long
                shape_info["paragraphs"] = []
                for p in shape.text_frame.paragraphs:
                    para_info = {"text": p.text, "level": p.level}
                    if p.runs:
                        run = p.runs[0]
                        para_info["font_size"] = str(run.font.size) if run.font.size else None
                        para_info["bold"] = run.font.bold
                    shape_info["paragraphs"].append(para_info)
            if shape.has_table:
                table = shape.table
                shape_info["table"] = {
                    "rows": len(table.rows),
                    "cols": len(table.columns),
                    "cells_preview": [
                        [table.cell(r, c).text[:50] for c in range(min(len(table.columns), 5))]
                        for r in range(min(len(table.rows), 5))
                    ],
                }
            slide_info["shapes"].append(shape_info)
        structure["slides"].append(slide_info)

    return json.dumps(structure, ensure_ascii=False, indent=2)


def inspect_slide_xml(pptx_bytes: bytes, slide_index: int) -> str:
    """Retourne le XML brut d'un slide."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    if slide_index >= len(prs.slides):
        return f"Erreur : slide {slide_index} n'existe pas (max: {len(prs.slides) - 1})"
    slide = prs.slides[slide_index]
    return etree.tostring(slide._element, pretty_print=True).decode()


# ============================================================
# Unpack / Repack PPTX (workflow d'édition XML)
# ============================================================

def unpack_pptx(pptx_bytes: bytes, dest_dir: str) -> str:
    """Décompresse un PPTX avec pretty-print XML et smart quotes."""
    unpacked_dir = str(Path(dest_dir) / "unpacked")
    return pptx_tools.unpack(pptx_bytes, unpacked_dir)


def repack_pptx(unpacked_dir: str, original_bytes: bytes = None) -> bytes:
    """
    Repackage avec validation complète, auto-repair, condensation XML et smart quotes.

    Stratégie (comme Claude le fait manuellement) :
    - Erreurs XSD sur slides → on les signale (le caller peut retenter)
    - Erreurs structurelles → on les signale (le caller peut retenter)
    - L'idée : ne jamais envoyer un fichier cassé, toujours retenter d'abord
    """
    # Nettoyer les fichiers orphelins
    pptx_tools.clean(unpacked_dir)

    # Validation complète (structurelle + XSD)
    validation_result = None
    try:
        validation_result = pptx_validate.validate_pptx(unpacked_dir, original_bytes)
        if validation_result["repairs"]:
            logger.info(f"Auto-repaired {validation_result['repairs']} issue(s)")
        if validation_result["valid"]:
            logger.info("Validation PPTX: OK")
    except Exception as e:
        logger.warning(f"Validation PPTX skippée (erreur): {e}")

    # Collecter les erreurs non-réparables
    blocking_errors = []
    if validation_result:
        # Erreurs structurelles = le fichier ne s'ouvrira probablement pas
        if validation_result["errors"]:
            blocking_errors.extend(
                f"[STRUCTURE] {e}" for e in validation_result["errors"]
            )
        # Erreurs XSD nouvelles = le LLM a introduit des tags/attributs invalides
        if validation_result["xsd_errors"]:
            blocking_errors.extend(
                f"[XSD] {e}" for e in validation_result["xsd_errors"]
            )

    if blocking_errors:
        raise ValueError(
            f"PPTX invalide — {len(blocking_errors)} erreur(s) détectée(s) :\n"
            + "\n".join(f"  • {e}" for e in blocking_errors[:5])
            + ("\n  ..." if len(blocking_errors) > 5 else "")
        )

    return pptx_tools.pack(unpacked_dir, original_bytes)


# ============================================================
# Appel LLM
# ============================================================

async def call_llm(system_prompt: str, query: str) -> str:
    """
    Appelle SiaGPT /plain_llm endpoint.
    Format : { systemPrompt, query, llm, temperature } → string
    """
    async with httpx.AsyncClient(timeout=120.0) as client:
        response = await client.post(
            LLM_API_URL,
            json={
                "systemPrompt": system_prompt,
                "query": query,
                "llm": LLM_MODEL,
                "temperature": 0.1,
            },
            headers={
                "Authorization": f"Bearer {LLM_API_KEY}",
                "Content-Type": "application/json",
            },
        )
        response.raise_for_status()

        # /plain_llm retourne directement un string
        data = response.json()
        if isinstance(data, str):
            return data
        # Au cas où c'est wrappé dans un objet
        if isinstance(data, dict):
            return data.get("content", data.get("text", str(data)))
        return str(data)


def extract_json(llm_response: str) -> dict:
    """Extrait le JSON de la réponse LLM (enlève les ```json si présents)."""
    text = llm_response.strip()
    # Enlever les blocs markdown
    if text.startswith("```json"):
        text = text[len("```json"):].strip()
    if text.startswith("```"):
        text = text[3:].strip()
    if text.endswith("```"):
        text = text[:-3].strip()
    return json.loads(text)


def extract_xml(llm_response: str) -> str:
    """Extrait le XML de la réponse LLM (enlève les ```xml si présents)."""
    text = llm_response.strip()
    if text.startswith("```xml"):
        text = text[len("```xml"):].strip()
    if text.startswith("```"):
        text = text[3:].strip()
    if text.endswith("```"):
        text = text[:-3].strip()
    return text


def validate_xml(xml_string: str) -> tuple[bool, str]:
    """Valide que le XML est bien formé. Retourne (valide, message_erreur)."""
    try:
        etree.fromstring(xml_string.encode("utf-8"))
        return True, ""
    except etree.XMLSyntaxError as e:
        return False, str(e)


def read_slide_xmls(unpacked_dir: str) -> dict[str, str]:
    """Lit tous les XML de slides depuis le dossier décompressé."""
    slides_dir = Path(unpacked_dir) / "ppt" / "slides"
    slides = {}
    if slides_dir.exists():
        for slide_file in sorted(slides_dir.glob("slide*.xml")):
            slides[slide_file.name] = slide_file.read_text(encoding="utf-8")
    return slides


async def plan_modifications(structure: str, prompt: str, slide_xmls: dict[str, str] = None) -> dict:
    """
    Phase 1 : Appelle le LLM pour planifier les modifications.
    Retourne un dict avec slides_to_modify, slides_to_add, slides_to_remove, summary.
    """
    query = (
        "PHASE : PLANIFICATION\n\n"
        f"Structure du fichier PPTX :\n{structure}\n\n"
    )

    # Ajouter un aperçu du contenu des slides si disponible
    if slide_xmls:
        query += "Contenu des slides (aperçu texte) :\n"
        for name, xml in slide_xmls.items():
            # Extraire juste le texte visible pour le planning
            texts = re.findall(r'<a:t[^>]*>([^<]+)</a:t>', xml)
            preview = " | ".join(texts[:20])  # Limiter l'aperçu
            query += f"  {name}: {preview[:300]}\n"
        query += "\n"

    query += (
        f"Demande de l'utilisateur : {prompt}\n\n"
        "Retourne UNIQUEMENT un JSON valide décrivant le plan de modifications."
    )

    for attempt in range(MAX_RETRIES):
        llm_response = await call_llm(SYSTEM_PROMPT, query)
        try:
            plan = extract_json(llm_response)
            # Valider la structure minimale
            if "summary" not in plan:
                plan["summary"] = "Modifications planifiées"
            return plan
        except (json.JSONDecodeError, ValueError) as e:
            if attempt < MAX_RETRIES - 1:
                query = (
                    f"Ta réponse précédente n'était pas du JSON valide.\n"
                    f"Erreur : {e}\n"
                    f"Ta réponse était :\n{llm_response[:500]}\n\n"
                    f"Retourne UNIQUEMENT un JSON valide. Pas de texte, pas de markdown."
                )
            else:
                raise ValueError(f"Le LLM n'a pas retourné de JSON valide après {MAX_RETRIES} tentatives")


async def modify_slide_xml(
    slide_xml: str,
    instructions: str,
    slide_name: str,
    structure_context: str = "",
) -> str:
    """
    Phase 2 : Appelle le LLM pour modifier le XML d'une slide.
    Retourne le XML modifié complet.
    """
    query = (
        "PHASE : MODIFICATION XML\n\n"
        f"Slide : {slide_name}\n\n"
    )
    if structure_context:
        query += f"Contexte de la présentation :\n{structure_context}\n\n"

    query += (
        f"Instructions : {instructions}\n\n"
        f"XML actuel de la slide :\n{slide_xml}\n\n"
        "Retourne UNIQUEMENT le XML modifié complet. Pas de markdown, pas d'explication."
    )

    for attempt in range(MAX_RETRIES):
        llm_response = await call_llm(SYSTEM_PROMPT, query)
        new_xml = extract_xml(llm_response)

        # Validation forte : parsing + XSD (détecte les tags inventés)
        is_valid, error_msg = pptx_validate.validate_slide_xml_string(new_xml)
        if is_valid:
            return new_xml

        # XML invalide → demander correction
        if attempt < MAX_RETRIES - 1:
            query = (
                "PHASE : MODIFICATION XML\n\n"
                f"Slide : {slide_name}\n\n"
                f"Ton XML précédent contenait une erreur : {error_msg}\n\n"
                f"XML que tu as retourné (début) :\n{new_xml[:1000]}\n\n"
                f"XML original de la slide :\n{slide_xml}\n\n"
                f"Instructions originales : {instructions}\n\n"
                "Corrige et retourne UNIQUEMENT le XML modifié complet et valide."
            )
        else:
            raise ValueError(f"XML invalide après {MAX_RETRIES} tentatives : {error_msg}")

    return slide_xml  # Fallback : retourner l'original


async def apply_xml_modifications(
    unpacked_dir: str,
    structure: str,
    prompt: str,
) -> dict:
    """
    Workflow complet XML pur :
    1. Lire les slides
    2. Planifier les modifications
    3. Appliquer les modifications XML slide par slide
    4. Retourne un résumé
    """
    slide_xmls = read_slide_xmls(unpacked_dir)
    slides_dir = Path(unpacked_dir) / "ppt" / "slides"

    # Phase 1 : Planifier
    plan = await plan_modifications(structure, prompt, slide_xmls)

    results = {
        "plan": plan,
        "modified_slides": [],
        "added_slides": [],
        "removed_slides": [],
        "errors": [],
    }

    # Phase 2a : Modifier les slides existantes
    for mod in plan.get("slides_to_modify", []):
        filename = mod["filename"]
        instructions = mod["instructions"]

        if filename not in slide_xmls:
            results["errors"].append(f"Slide {filename} introuvable")
            continue

        try:
            new_xml = await modify_slide_xml(
                slide_xmls[filename],
                instructions,
                filename,
                structure_context=plan.get("summary", ""),
            )
            # Écrire le XML modifié
            (slides_dir / filename).write_text(new_xml, encoding="utf-8")
            results["modified_slides"].append(filename)
        except Exception as e:
            results["errors"].append(f"Erreur sur {filename}: {str(e)}")

    # Phase 2b : Ajouter des slides (duplication + modification)
    for add in plan.get("slides_to_add", []):
        source = add.get("duplicate_from", "")
        instructions = add.get("instructions", "")
        position = add.get("position", None)

        if source not in slide_xmls:
            results["errors"].append(f"Slide source {source} introuvable pour duplication")
            continue

        try:
            # Dupliquer la slide via pptx_tools (gère .rels, Content_Types, notesSlide)
            dup_info = pptx_tools.duplicate_slide(unpacked_dir, source)
            new_filename = dup_info["new_filename"]

            # Ajouter dans presentation.xml à la bonne position
            pptx_tools.add_slide_to_presentation(
                unpacked_dir,
                dup_info["new_sld_id"],
                dup_info["new_r_id"],
                position=position,
            )

            # Modifier le contenu si des instructions sont fournies
            if instructions:
                new_slide_xml = (slides_dir / new_filename).read_text(encoding="utf-8")
                modified_xml = await modify_slide_xml(
                    new_slide_xml,
                    instructions,
                    new_filename,
                    structure_context=plan.get("summary", ""),
                )
                (slides_dir / new_filename).write_text(modified_xml, encoding="utf-8")

            results["added_slides"].append(new_filename)
        except Exception as e:
            results["errors"].append(f"Erreur ajout slide depuis {source}: {str(e)}")

    # Phase 2c : Supprimer des slides
    # On retire juste le <p:sldId> de presentation.xml
    # Le nettoyage des fichiers orphelins est fait par clean() au moment du repack
    for filename in plan.get("slides_to_remove", []):
        try:
            pres_path = Path(unpacked_dir) / "ppt" / "presentation.xml"
            pres_xml = pres_path.read_text(encoding="utf-8")

            # Trouver le rId correspondant au fichier
            pres_rels_path = Path(unpacked_dir) / "ppt" / "_rels" / "presentation.xml.rels"
            pres_rels = pres_rels_path.read_text(encoding="utf-8")
            r_id_match = re.search(rf'Id="(rId\d+)"[^>]*Target="slides/{filename}"', pres_rels)

            if r_id_match:
                r_id = r_id_match.group(1)
                # Retirer le sldId de presentation.xml
                pres_xml = re.sub(rf'\s*<p:sldId[^>]*r:id="{r_id}"[^>]*/>', '', pres_xml)
                pres_path.write_text(pres_xml, encoding="utf-8")

                results["removed_slides"].append(filename)
            else:
                results["errors"].append(f"Slide {filename} non trouvée dans les relations")
        except Exception as e:
            results["errors"].append(f"Erreur suppression {filename}: {str(e)}")

    return results


# ============================================================
# Fonctions core — logique partagée REST / MCP
# ============================================================

async def _do_edit(pptx_bytes: bytes, prompt: str, auth_token: str, output_filename: str = None) -> dict:
    """Logique core d'édition PPTX. Utilisée par REST et MCP."""
    if not output_filename:
        output_filename = f"modified_{uuid.uuid4().hex[:8]}.pptx"

    structure = inspect_pptx_structure(pptx_bytes)

    with tempfile.TemporaryDirectory() as tmp_dir:
        unpacked_dir = unpack_pptx(pptx_bytes, tmp_dir)
        results = await apply_xml_modifications(unpacked_dir, structure, prompt)
        output_bytes = repack_pptx(unpacked_dir, pptx_bytes)

    media_info = await save_to_siagpt_medias(output_bytes, output_filename, auth_token)

    return {
        "status": "ok",
        "summary": results["plan"].get("summary", ""),
        "modified_slides": results["modified_slides"],
        "added_slides": results["added_slides"],
        "removed_slides": results["removed_slides"],
        "errors": results["errors"],
        "media_uuid": media_info.get("uuid"),
        "media_name": media_info.get("name"),
    }


async def _do_create(prompt: str, auth_token: str, template_bytes: bytes = None, output_filename: str = None) -> dict:
    """Logique core de création PPTX. Utilisée par REST et MCP."""
    if not output_filename:
        output_filename = f"new_{uuid.uuid4().hex[:8]}.pptx"

    if not template_bytes:
        template_bytes = create_skeleton_pptx(prompt)

    create_prompt = (
        f"CRÉATION DE PRÉSENTATION depuis un template.\n\n"
        f"Demande : {prompt}\n\n"
        f"Modifie les slides existantes pour répondre à la demande. "
        f"Tu peux dupliquer des slides pour en ajouter, en supprimer si nécessaire, "
        f"et modifier tout le contenu texte."
    )

    return await _do_edit(template_bytes, create_prompt, auth_token, output_filename)


def _format_mcp_summary(action: str, result: dict, extra_line: str = None) -> str:
    """Formate un résumé texte pour les réponses MCP."""
    lines = [f"Présentation {action} avec succès !"]
    if extra_line:
        lines.append(f"- {extra_line}")
    lines.append(f"- Fichier : {result.get('media_name', 'N/A')}")
    lines.append(f"- UUID : {result.get('media_uuid', 'N/A')}")
    lines.append(f"- {result.get('summary', '')}")
    if result.get("modified_slides"):
        lines.append(f"- Slides modifiées : {', '.join(result['modified_slides'])}")
    if result.get("added_slides"):
        lines.append(f"- Slides ajoutées : {', '.join(result['added_slides'])}")
    if result.get("errors"):
        lines.append(f"- Avertissements : {'; '.join(result['errors'])}")
    return "\n".join(lines)


# ============================================================
# Endpoint principal — Modification de PPTX
# ============================================================

@app.post("/api/edit")
async def edit_pptx(
    request: Request,
    prompt: str = Form(...),
    file: UploadFile = File(...),
    output_filename: str = Form(None),
):
    """Modifie un PPTX existant. Mode XML pur."""
    auth_token = (request.headers.get("authorization", "").removeprefix("Bearer ").strip()) or LLM_API_KEY
    pptx_bytes = await file.read()
    try:
        return await _do_edit(pptx_bytes, prompt, auth_token, output_filename)
    except ValueError as e:
        raise HTTPException(status_code=500, detail=str(e))


# ============================================================
# Endpoint — Création de PPTX
# ============================================================

@app.post("/api/create")
async def create_pptx(
    request: Request,
    prompt: str = Form(...),
    template: UploadFile = File(None),
    output_filename: str = Form(None),
):
    """Crée un PPTX depuis un template (ou un squelette vierge). Mode XML pur."""
    auth_token = (request.headers.get("authorization", "").removeprefix("Bearer ").strip()) or LLM_API_KEY
    template_bytes = await template.read() if template else None
    try:
        return await _do_create(prompt, auth_token, template_bytes, output_filename)
    except ValueError as e:
        raise HTTPException(status_code=500, detail=str(e))


def create_skeleton_pptx(prompt: str) -> bytes:
    """
    Crée un PPTX squelette basique quand aucun template n'est fourni.
    C'est du code contrôlé (pas du LLM), donc pas de risque sécu.
    """
    prs = Presentation()
    # Créer quelques slides vierges avec des placeholders
    # Le LLM les remplira ensuite via XML
    for i in range(5):
        slide_layout = prs.slide_layouts[5]  # Layout "Blank"
        slide = prs.slides.add_slide(slide_layout)
        # Ajouter un textbox titre
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tf = txBox.text_frame
        tf.text = f"[Titre slide {i+1}]"
        # Ajouter un textbox contenu
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf2 = txBox2.text_frame
        tf2.text = f"[Contenu slide {i+1}]"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ============================================================
# Endpoint — Inspection
# ============================================================

@app.post("/api/inspect")
async def inspect_pptx(file: UploadFile = File(...)):
    """Retourne la structure d'un PPTX en JSON."""
    pptx_bytes = await file.read()
    structure = inspect_pptx_structure(pptx_bytes)
    return JSONResponse(content=json.loads(structure))


@app.post("/api/inspect/xml")
async def inspect_xml(file: UploadFile = File(...), slide_index: int = Form(0)):
    """Retourne le XML brut d'un slide."""
    pptx_bytes = await file.read()
    xml = inspect_slide_xml(pptx_bytes, slide_index)
    return {"slide_index": slide_index, "xml": xml}



# ============================================================
# MCP Server — SSE Transport
# ============================================================

# Sessions MCP actives : session_id → asyncio.Queue
mcp_sessions: dict[str, asyncio.Queue] = {}


def mcp_jsonrpc_response(req_id, result):
    """Construit une réponse JSON-RPC 2.0."""
    return {"jsonrpc": "2.0", "id": req_id, "result": result}


def mcp_jsonrpc_error(req_id, code, message):
    """Construit une erreur JSON-RPC 2.0."""
    return {"jsonrpc": "2.0", "id": req_id, "error": {"code": code, "message": message}}


@app.get("/mcp/sse")
async def mcp_sse_get(request: Request):
    """
    Endpoint SSE pour le protocole MCP (ancien transport).
    """
    session_id = uuid.uuid4().hex
    queue: asyncio.Queue = asyncio.Queue()
    mcp_sessions[session_id] = queue

    async def event_stream():
        scheme = request.headers.get("x-forwarded-proto", "https")
        host = request.headers.get("host", request.base_url.hostname)
        endpoint_url = f"{scheme}://{host}/mcp/messages?session_id={session_id}"
        yield f"event: endpoint\ndata: {endpoint_url}\n\n"

        try:
            while True:
                if await request.is_disconnected():
                    break
                try:
                    message = await asyncio.wait_for(queue.get(), timeout=30.0)
                    yield f"event: message\ndata: {json.dumps(message)}\n\n"
                except asyncio.TimeoutError:
                    yield ": keepalive\n\n"
        finally:
            mcp_sessions.pop(session_id, None)

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


async def handle_mcp_request(body: dict, session_id: str = "") -> tuple[dict, str]:
    """
    Traite une requête JSON-RPC MCP et retourne (réponse, session_id).
    """
    method = body.get("method", "")
    req_id = body.get("id")
    params = body.get("params", {})

    # --- initialize ---
    if method == "initialize":
        if not session_id:
            session_id = uuid.uuid4().hex
        return mcp_jsonrpc_response(req_id, {
            "protocolVersion": "2024-11-05",
            "capabilities": {"tools": {"listChanged": False}},
            "serverInfo": {"name": "pptx-service", "version": "1.0.0"},
        }), session_id

    # --- notifications/initialized ---
    if method == "notifications/initialized":
        return None, session_id

    # --- tools/list ---
    if method == "tools/list":
        return mcp_jsonrpc_response(req_id, {
            "tools": [
                {
                    "name": "generate_pptx",
                    "description": "Génère une présentation PowerPoint à partir d'une description textuelle. Peut utiliser un template existant comme base (recommandé pour les présentations Sia Partners). Le fichier est sauvegardé dans la collection SiaGPT.",
                    "inputSchema": {
                        "type": "object",
                        "properties": {
                            "prompt": {
                                "type": "string",
                                "description": "Description de la présentation à créer (contenu, nombre de slides, style...)",
                            },
                            "template_file_id": {
                                "type": "string",
                                "description": "UUID d'un template PPTX dans la collection SiaGPT à utiliser comme base. Si omis, crée un squelette vierge.",
                            }
                        },
                        "required": ["prompt"],
                    },
                },
                {
                    "name": "edit_pptx",
                    "description": "Modifie une présentation PowerPoint existante dans la collection SiaGPT. Récupère le fichier par son UUID, applique les modifications demandées, et uploade la version modifiée.",
                    "inputSchema": {
                        "type": "object",
                        "properties": {
                            "prompt": {
                                "type": "string",
                                "description": "Description des modifications à apporter (ex: changer les couleurs, ajouter une slide, modifier le texte...)",
                            },
                            "source_file_id": {
                                "type": "string",
                                "description": "UUID du fichier PPTX dans la collection SiaGPT à modifier",
                            }
                        },
                        "required": ["prompt", "source_file_id"],
                    },
                }
            ]
        }), session_id

    # --- tools/call ---
    if method == "tools/call":
        tool_name = params.get("name", "")
        tool_args = params.get("arguments", {})

        if tool_name == "generate_pptx":
            prompt = tool_args.get("prompt", "")
            template_file_id = tool_args.get("template_file_id", "")
            if not prompt:
                return mcp_jsonrpc_error(req_id, -32602, "Le paramètre 'prompt' est requis"), session_id

            try:
                # Si un template est fourni, le télécharger depuis SiaGPT Medias
                template_bytes = None
                template_info = ""
                if template_file_id:
                    template_bytes, template_name = await download_from_siagpt_medias(template_file_id, LLM_API_KEY)
                    template_info = f"Template : {template_name} ({template_file_id})"

                result = await _do_create(prompt, LLM_API_KEY, template_bytes)
                summary = _format_mcp_summary("créée", result, template_info)
                return mcp_jsonrpc_response(req_id, {
                    "content": [{"type": "text", "text": summary}]
                }), session_id
            except Exception as e:
                return mcp_jsonrpc_error(req_id, -32000, str(e)), session_id

        if tool_name == "edit_pptx":
            prompt = tool_args.get("prompt", "")
            source_file_id = tool_args.get("source_file_id", "")
            if not prompt:
                return mcp_jsonrpc_error(req_id, -32602, "Le paramètre 'prompt' est requis"), session_id
            if not source_file_id:
                return mcp_jsonrpc_error(req_id, -32602, "Le paramètre 'source_file_id' est requis"), session_id

            try:
                pptx_bytes, original_filename = await download_from_siagpt_medias(source_file_id, LLM_API_KEY)
                result = await _do_edit(pptx_bytes, prompt, LLM_API_KEY)
                summary = _format_mcp_summary("modifiée", result, f"Source : {original_filename} ({source_file_id})")
                return mcp_jsonrpc_response(req_id, {
                    "content": [{"type": "text", "text": summary}]
                }), session_id
            except httpx.HTTPStatusError as e:
                return mcp_jsonrpc_error(req_id, -32000, f"Fichier {source_file_id} introuvable : {e.response.status_code}"), session_id
            except Exception as e:
                return mcp_jsonrpc_error(req_id, -32000, str(e)), session_id

        return mcp_jsonrpc_error(req_id, -32601, f"Tool inconnu : {tool_name}"), session_id

    return mcp_jsonrpc_error(req_id, -32601, f"Méthode inconnue : {method}"), session_id


@app.post("/mcp/sse")
async def mcp_sse_post(request: Request):
    """
    Endpoint Streamable HTTP — POST direct avec réponse JSON-RPC.
    """
    body = await request.json()
    session_id = request.headers.get("mcp-session-id", "")
    response, session_id = await handle_mcp_request(body, session_id)
    
    if response is None:
        # Notification — pas de réponse body
        return JSONResponse(
            content={},
            status_code=202,
            headers={"mcp-session-id": session_id},
        )
    
    return JSONResponse(
        content=response,
        headers={"mcp-session-id": session_id},
    )


@app.delete("/mcp/sse")
async def mcp_sse_delete(request: Request):
    """Fermeture de session MCP."""
    return JSONResponse({"status": "ok"})


@app.get("/mcp/messages")
async def mcp_messages_get(request: Request, session_id: str = ""):
    """
    GET sur /mcp/messages — Langflow vérifie l'endpoint ou ouvre un stream SSE.
    """
    return JSONResponse({"status": "ok", "message": "Use POST to send MCP messages"})


@app.post("/mcp/messages")
async def mcp_messages(request: Request, session_id: str = ""):
    """
    Endpoint messages pour le transport SSE classique.
    """
    if session_id not in mcp_sessions:
        body = await request.json()
        response, _ = await handle_mcp_request(body, session_id)
        if response is None:
            return JSONResponse({}, status_code=202)
        return JSONResponse(response)

    queue = mcp_sessions[session_id]
    body = await request.json()
    response, _ = await handle_mcp_request(body, session_id)
    if response is not None:
        await queue.put(response)
    return JSONResponse({"status": "ok"})


# ============================================================
# Endpoint unifié — Génère ou modifie un PPTX automatiquement
# ============================================================

@app.post("/api/generate")
async def generate_pptx(request: Request):
    """Endpoint unifié — accepte JSON ou form-data, crée ou modifie un PPTX."""
    auth_token = (request.headers.get("authorization", "").removeprefix("Bearer ").strip()) or LLM_API_KEY
    content_type = request.headers.get("content-type", "")

    if "application/json" in content_type:
        body = await request.json()
        prompt = body.get("prompt", "")
        if not prompt:
            raise HTTPException(status_code=400, detail="Le champ 'prompt' est requis")

        # Si template_file_id fourni, télécharger le template depuis SiaGPT Medias
        template_bytes = None
        template_file_id = body.get("template_file_id", "")
        if template_file_id:
            template_bytes, _ = await download_from_siagpt_medias(template_file_id, auth_token)

        return await _do_create(prompt, auth_token, template_bytes, output_filename=body.get("output_filename"))
    else:
        form = await request.form()
        prompt = form.get("prompt", "")
        if not prompt:
            raise HTTPException(status_code=400, detail="Le champ 'prompt' est requis")
        file = form.get("file", None)
        output_filename = form.get("output_filename", None)
        if file and hasattr(file, 'filename') and file.filename:
            pptx_bytes = await file.read()
            return await _do_edit(pptx_bytes, prompt, auth_token, output_filename)
        else:
            # En form-data, template_file_id aussi supporté
            template_bytes = None
            template_file_id = form.get("template_file_id", "")
            if template_file_id:
                template_bytes, _ = await download_from_siagpt_medias(template_file_id, auth_token)
            return await _do_create(prompt, auth_token, template_bytes, output_filename=output_filename)


# ============================================================
# Health check
# ============================================================

@app.get("/health")
async def health():
    """Health check — retourne le statut et la configuration du service."""
    return {
        "status": "ok",
        "mode": "xml-pure",
        "exec_enabled": False,
        "llm_configured": bool(LLM_API_KEY),
        "collection_configured": bool(SIAGPT_COLLECTION_ID),
        "model": LLM_MODEL,
    }


@app.api_route("/", methods=["GET", "POST", "DELETE"])
async def root(request: Request):
    """Racine — healthcheck et fallback MCP."""
    if request.method == "POST":
        try:
            body = await request.json()
            if "jsonrpc" in body:
                # C'est une requête MCP
                session_id = request.headers.get("mcp-session-id", "")
                response, session_id = await handle_mcp_request(body, session_id)
                if response is None:
                    return JSONResponse({}, status_code=202, headers={"mcp-session-id": session_id})
                return JSONResponse(content=response, headers={"mcp-session-id": session_id})
        except Exception:
            pass
    return JSONResponse({"status": "ok", "service": "pptx-service"})
